"""Document ingestion - chunk, embed, and upsert per-log corpora.

Adapted from the original repo's `backend/utils/ingest_documents.py`. Key
differences from the original:

  * Inputs are passed in (no .env/global state): `documents_dir`, `namespace`,
    `index`, `embedder`, and a pre-built vision `llm` for image / pptx
    extraction. Provider/model selection happens upstream.
  * Targets a per-log Pinecone namespace (`cde-{log_id}`) so corpora from
    different logs don't collide.
  * `unstructured` is dropped in favour of `pypdf` + `python-docx` +
    `python-pptx` to keep the dependency tree light. Plain `.txt` files are
    read directly.
"""

from __future__ import annotations

import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Callable, Iterable, Optional

import pypdf
from langchain_core.language_models import BaseChatModel
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

from .image_analyzer import analyze_image_content

KB_NS = "bpm-kb"
SUPPORTED_SUFFIXES = {".pdf", ".pptx", ".png", ".jpg", ".jpeg", ".docx", ".txt"}


def get_timestamp_from_filename(filename: str) -> Optional[int]:
    """Parse a leading `YYYY-MM-DD_` prefix into a Unix timestamp."""
    try:
        date_str = filename.split("_", 1)[0]
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return int(dt.timestamp())
    except (ValueError, IndexError):
        return None


def _load_pdf(path: Path) -> list[str]:
    reader = pypdf.PdfReader(str(path))
    return [page.extract_text() or "" for page in reader.pages]


def _load_docx(path: Path) -> list[str]:
    # Imported lazily; python-docx is a relatively heavy import.
    from docx import Document

    doc = Document(str(path))
    return ["\n".join(p.text for p in doc.paragraphs if p.text.strip())]


def _load_txt(path: Path) -> list[str]:
    return [path.read_text(encoding="utf-8", errors="ignore")]


def _load_pptx(path: Path, *, vision_llm: BaseChatModel, cache_dir: Path) -> list[str]:
    prs = Presentation(str(path))
    texts: list[str] = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                tmp = cache_dir / f"temp_{image.sha1}.{image.ext}"
                try:
                    tmp.write_bytes(image.blob)
                    description = analyze_image_content(tmp, llm=vision_llm)
                    if "Error" not in description:
                        texts.append(
                            f"Description of an image from slide {i + 1} of "
                            f"'{path.name}': {description}"
                        )
                finally:
                    if tmp.exists():
                        tmp.unlink()
    return texts


def _load_image(path: Path, *, vision_llm: BaseChatModel) -> list[str]:
    description = analyze_image_content(path, llm=vision_llm)
    if "Error" in description:
        return []
    return [description]


def _extract_texts(
    path: Path,
    *,
    vision_llm: BaseChatModel,
    cache_dir: Path,
) -> list[str]:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _load_pdf(path)
        if suffix == ".docx":
            return _load_docx(path)
        if suffix == ".txt":
            return _load_txt(path)
        if suffix == ".pptx":
            return _load_pptx(path, vision_llm=vision_llm, cache_dir=cache_dir)
        if suffix in {".png", ".jpg", ".jpeg"}:
            return _load_image(path, vision_llm=vision_llm)
    except Exception as e:
        logging.error("Failed to read %s: %s", path.name, e)
        return []
    return []


def process_and_embed(
    *,
    index,
    embedder,
    text_splitter: RecursiveCharacterTextSplitter,
    texts_to_embed: list[str],
    source_document_name: str,
    doc_timestamp: int,
    namespace: str,
) -> int:
    if not texts_to_embed:
        return 0

    documents = text_splitter.create_documents(texts_to_embed)
    if not documents:
        return 0

    texts = [chunk.page_content for chunk in documents]
    vectors = embedder.embed_documents(texts)

    vectors_to_upsert = []
    stem = Path(source_document_name).stem
    for i, chunk in enumerate(documents):
        vector_id = f"{stem}_{abs(hash(chunk.page_content))}_{i}"
        metadata = {
            "text": chunk.page_content,
            "source": source_document_name,
            "timestamp": doc_timestamp,
        }
        vectors_to_upsert.append((vector_id, vectors[i], metadata))

    index.upsert(vectors=vectors_to_upsert, namespace=namespace)
    return len(vectors_to_upsert)


def process_context_files(
    files: Iterable[Path],
    *,
    index,
    embedder,
    vision_llm: BaseChatModel,
    namespace: str,
    cache_dir: Optional[Path] = None,
    progress: Optional[Callable[[float, str], None]] = None,
) -> dict:
    """Ingest a batch of files into the given Pinecone namespace.

    Returns a summary `{ ingested: N, skipped: [...], failed: [...] }`.
    """
    files = list(files)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    cache_dir = cache_dir or Path(os.path.join(os.getcwd(), ".cde-tmp"))
    cache_dir.mkdir(parents=True, exist_ok=True)

    total_vectors = 0
    skipped: list[dict] = []
    failed: list[dict] = []

    for i, path in enumerate(files):
        if progress:
            progress(i / max(1, len(files)), f"Indexing {path.name}")

        ts = get_timestamp_from_filename(path.name)
        if ts is None:
            skipped.append({"file": path.name, "reason": "missing YYYY-MM-DD_ prefix"})
            continue

        texts = _extract_texts(path, vision_llm=vision_llm, cache_dir=cache_dir)
        if not texts:
            failed.append({"file": path.name, "reason": "no text extracted"})
            continue

        try:
            count = process_and_embed(
                index=index,
                embedder=embedder,
                text_splitter=text_splitter,
                texts_to_embed=texts,
                source_document_name=path.name,
                doc_timestamp=ts,
                namespace=namespace,
            )
            total_vectors += count
        except Exception as e:
            failed.append({"file": path.name, "reason": str(e)})

    if progress:
        progress(1.0, "Indexing complete")

    return {"ingested": total_vectors, "skipped": skipped, "failed": failed}


def process_glossary_file(
    glossary_path: Path,
    *,
    index,
    embedder,
    namespace: str = KB_NS,
) -> int:
    """Embed the BPM glossary into the `bpm-kb` namespace."""
    import pandas as pd  # local - pandas is inherited from the platform

    if not glossary_path.exists():
        return 0

    df = pd.read_csv(glossary_path)
    if "term" not in df.columns or "definition" not in df.columns:
        return 0

    df["text_to_embed"] = df["term"].astype(str) + ": " + df["definition"].astype(str)
    texts = df["text_to_embed"].tolist()
    vectors = embedder.embed_documents(texts)

    vectors_to_upsert = []
    for i, row in df.iterrows():
        vectors_to_upsert.append(
            (
                f"bpm_kb_{i}",
                vectors[i],
                {"text": row["text_to_embed"], "source": "BPM Glossary", "timestamp": 0},
            )
        )

    index.upsert(vectors=vectors_to_upsert, namespace=namespace)
    return len(vectors_to_upsert)


def delete_document_vectors(index, *, namespace: str, source_document_name: str) -> None:
    """Remove all chunks for a given source document from the namespace."""
    try:
        index.delete(filter={"source": source_document_name}, namespace=namespace)
    except Exception as e:
        logging.error("Failed to delete vectors for %s: %s", source_document_name, e)
